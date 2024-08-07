import collections.abc
import requests
import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import openai
from openai import OpenAI
from io import BytesIO
import config

# API Token
openai.api_key = config.API_KEY
client = OpenAI(api_key=openai.api_key)

def slide_generator(text, prs, num_images=1, words_per_slide=200):
    words = text.split()
    slide_text = ' '.join(words[:words_per_slide])

    # Generate DALL-E prompt and get image
    prompt = f"Summarize the following text to a DALL-E image generation prompt: \n {slide_text}"
    dlp = client.chat.completions.create(
        model="gpt-3.5-turbo",
        messages=[{"role": "user", "content": f"{prompt}"}],
        max_tokens=250,
        n=1,
        stop=None,
        temperature=0.8
    )
    dalle_prompt = dlp.choices[0].message.content.strip()

    image_urls = []
    for _ in range(num_images):
        response = client.images.generate(
            model="dall-e-3",
            prompt=dalle_prompt + " Style: digital art",
            quality="standard",
            n=1,
            size="1024x1024"
        )
        image_urls.append(response.data[0].url)

    # Generate bullet point text for slide
    prompt = f"Create a bullet point text for a PowerPoint slide from the following text: \n {slide_text}"
    ppt = client.chat.completions.create(
        model="gpt-3.5-turbo",
        messages=[{"role": "user", "content": f"{prompt}"}],
        max_tokens=1024,
        n=1,
        stop=None,
        temperature=0.8
    )
    ppt_text = ppt.choices[0].message.content.strip()

    # Generate title for slide
    prompt = f"Create a title for a PowerPoint slide from the following text: \n {slide_text}"
    ppt = client.chat.completions.create(
        model="gpt-3.5-turbo",
        messages=[{"role": "user", "content": f"{prompt}"}],
        max_tokens=1024,
        n=1,
        stop=None,
        temperature=0.8
    )
    ppt_header = ppt.choices[0].message.content.strip()

    # Add a blank slide
    slide_layout = prs.slide_layouts[6]  # 6 is a completely blank slide
    slide = prs.slides.add_slide(slide_layout)

    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), prs.slide_width - Inches(1), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = ppt_header
    title_frame.paragraphs[0].font.name = 'Calibri (Body)'
    title_frame.paragraphs[0].font.size = Pt(28)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Add image
    image_y_position = Inches(1.5)
    image_height = prs.slide_height * 0.33
    for image_url in image_urls:
        response = requests.get(image_url)
        img_bytes = BytesIO(response.content)
        slide.shapes.add_picture(img_bytes, Inches(0), image_y_position, width=prs.slide_width, height=image_height)
        image_y_position += image_height

    # Add text box for slide content
    text_y_position = image_y_position + Inches(0.5)
    text_height = prs.slide_height * 0.33
    text_box = slide.shapes.add_textbox(Inches(0.5), text_y_position, prs.slide_width - Inches(1), text_height)
    text_frame = text_box.text_frame
    text_frame.text = ppt_text
    for paragraph in text_frame.paragraphs:
        paragraph.font.name = 'Calibri (Body)'
        paragraph.font.size = Pt(18)

# Streamlit App
st.title("Create PPT Slides with OpenAI")

num_slides = st.number_input("Enter the number of slides:", min_value=1, step=1)
prompts = []

if 'current_slide' not in st.session_state:
    st.session_state.current_slide = 0

if 'prs' not in st.session_state:
    st.session_state.prs = Presentation()
    width = Inches(13.33)
    height = Inches(7.5)
    st.session_state.prs.slide_width = width
    st.session_state.prs.slide_height = height

for i in range(num_slides):
    prompt = st.text_input(f"Enter prompt for slide {i+1}:")
    prompts.append(prompt)

if st.button("Create Slides"):
    for idx, prompt in enumerate(prompts):
        if prompt:
            st.write(f"Processing slide {idx+1}")
            slide_generator(prompt, st.session_state.prs)
            st.write(f"Slide {idx+1} created. Sending the next prompt.")
    st.session_state.prs.save("my_presentation.pptx")
    st.success("Presentation created successfully!")
    with open("my_presentation.pptx", "rb") as f:
        st.download_button(
            label="Download PPT",
            data=f,
            file_name="my_presentation.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
